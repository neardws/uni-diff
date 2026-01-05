#!/usr/bin/env python3
"""Generate binary test fixtures (images, Office docs, PDF)."""

import os
import sys

FIXTURES_DIR = os.path.dirname(os.path.abspath(__file__))


def generate_images():
    """Generate test image files."""
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        print("Pillow not installed, skipping image generation")
        return False

    images_dir = os.path.join(FIXTURES_DIR, 'fixtures', 'images')
    os.makedirs(images_dir, exist_ok=True)

    def create_image(text, filename, bg_color, text_color):
        img = Image.new('RGB', (400, 200), color=bg_color)
        draw = ImageDraw.Draw(img)
        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 24)
        except (IOError, OSError):
            font = ImageFont.load_default()
        draw.text((50, 80), text, fill=text_color, font=font)
        img.save(os.path.join(images_dir, filename))

    create_image("Hello World\nVersion 1.0", "old.png", "white", "black")
    create_image("Hello Universe\nVersion 2.0", "new.png", "white", "black")

    create_image("Original Image\nTest Content", "old.jpg", "lightblue", "darkblue")
    create_image("Modified Image\nUpdated Content", "new.jpg", "lightgreen", "darkgreen")

    old_gif = Image.new('RGB', (200, 100), color='yellow')
    draw = ImageDraw.Draw(old_gif)
    draw.text((20, 40), "Old GIF", fill='black')
    old_gif.save(os.path.join(images_dir, 'old.gif'))

    new_gif = Image.new('RGB', (200, 100), color='orange')
    draw = ImageDraw.Draw(new_gif)
    draw.text((20, 40), "New GIF", fill='black')
    new_gif.save(os.path.join(images_dir, 'new.gif'))

    print("Images generated successfully")
    return True


def generate_word_docs():
    """Generate test Word documents."""
    try:
        from docx import Document
        from docx.shared import Inches
    except ImportError:
        print("python-docx not installed, skipping Word generation")
        return False

    office_dir = os.path.join(FIXTURES_DIR, 'fixtures', 'office')
    os.makedirs(office_dir, exist_ok=True)

    doc = Document()
    doc.add_heading('Document Title', 0)
    doc.add_paragraph('This is the first paragraph of the document.')
    doc.add_paragraph('Second paragraph with some content.')
    doc.add_heading('Section 1', level=1)
    doc.add_paragraph('Content under section 1.')
    doc.save(os.path.join(office_dir, 'old.docx'))

    doc = Document()
    doc.add_heading('Document Title Updated', 0)
    doc.add_paragraph('This is the modified first paragraph.')
    doc.add_paragraph('Second paragraph with updated content.')
    doc.add_paragraph('A new third paragraph added.')
    doc.add_heading('Section 1', level=1)
    doc.add_paragraph('Updated content under section 1.')
    doc.add_heading('Section 2', level=1)
    doc.add_paragraph('New section with new content.')
    doc.save(os.path.join(office_dir, 'new.docx'))

    print("Word documents generated successfully")
    return True


def generate_excel_files():
    """Generate test Excel files."""
    try:
        from openpyxl import Workbook
    except ImportError:
        print("openpyxl not installed, skipping Excel generation")
        return False

    office_dir = os.path.join(FIXTURES_DIR, 'fixtures', 'office')
    os.makedirs(office_dir, exist_ok=True)

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws['A1'] = 'Name'
    ws['B1'] = 'Value'
    ws['A2'] = 'Item 1'
    ws['B2'] = 100
    ws['A3'] = 'Item 2'
    ws['B3'] = 200
    ws['A4'] = 'Item 3'
    ws['B4'] = 300
    wb.save(os.path.join(office_dir, 'old.xlsx'))

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws['A1'] = 'Name'
    ws['B1'] = 'Value'
    ws['C1'] = 'Category'
    ws['A2'] = 'Item 1'
    ws['B2'] = 150
    ws['C2'] = 'A'
    ws['A3'] = 'Item 2'
    ws['B3'] = 200
    ws['C3'] = 'B'
    ws['A4'] = 'Item 3'
    ws['B4'] = 350
    ws['C4'] = 'A'
    ws['A5'] = 'Item 4'
    ws['B5'] = 400
    ws['C5'] = 'C'
    wb.save(os.path.join(office_dir, 'new.xlsx'))

    print("Excel files generated successfully")
    return True


def generate_powerpoint_files():
    """Generate test PowerPoint files."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("python-pptx not installed, skipping PowerPoint generation")
        return False

    office_dir = os.path.join(FIXTURES_DIR, 'fixtures', 'office')
    os.makedirs(office_dir, exist_ok=True)

    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Presentation Title"
    subtitle.text = "Version 1.0"

    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "First Slide"
    body.text = "Content of first slide"
    prs.save(os.path.join(office_dir, 'old.pptx'))

    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Presentation Title Updated"
    subtitle.text = "Version 2.0"

    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "First Slide Modified"
    body.text = "Updated content of first slide"

    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "New Second Slide"
    body.text = "Content of new slide"
    prs.save(os.path.join(office_dir, 'new.pptx'))

    print("PowerPoint files generated successfully")
    return True


def generate_pdf_files():
    """Generate test PDF files."""
    try:
        from fpdf import FPDF
    except ImportError:
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            return generate_pdf_with_reportlab()
        except ImportError:
            print("Neither fpdf2 nor reportlab installed, skipping PDF generation")
            return False

    pdf_dir = os.path.join(FIXTURES_DIR, 'fixtures', 'pdf')
    os.makedirs(pdf_dir, exist_ok=True)

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=16)
    pdf.cell(200, 10, txt="Document Title", ln=True, align='C')
    pdf.set_font("Arial", size=12)
    pdf.multi_cell(0, 10, txt="This is the first paragraph of the PDF document.\n\nSecond paragraph with some content.\n\nThird paragraph here.")
    pdf.output(os.path.join(pdf_dir, 'old.pdf'))

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=16)
    pdf.cell(200, 10, txt="Document Title - Updated", ln=True, align='C')
    pdf.set_font("Arial", size=12)
    pdf.multi_cell(0, 10, txt="This is the modified first paragraph.\n\nSecond paragraph with updated content.\n\nA new third paragraph added.\n\nFourth paragraph with more information.")
    pdf.output(os.path.join(pdf_dir, 'new.pdf'))

    print("PDF files generated successfully")
    return True


def generate_pdf_with_reportlab():
    """Fallback PDF generation using reportlab."""
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    pdf_dir = os.path.join(FIXTURES_DIR, 'fixtures', 'pdf')
    os.makedirs(pdf_dir, exist_ok=True)

    c = canvas.Canvas(os.path.join(pdf_dir, 'old.pdf'), pagesize=letter)
    c.setFont("Helvetica-Bold", 16)
    c.drawString(200, 750, "Document Title")
    c.setFont("Helvetica", 12)
    c.drawString(50, 700, "This is the first paragraph of the PDF document.")
    c.drawString(50, 680, "Second paragraph with some content.")
    c.drawString(50, 660, "Third paragraph here.")
    c.save()

    c = canvas.Canvas(os.path.join(pdf_dir, 'new.pdf'), pagesize=letter)
    c.setFont("Helvetica-Bold", 16)
    c.drawString(200, 750, "Document Title - Updated")
    c.setFont("Helvetica", 12)
    c.drawString(50, 700, "This is the modified first paragraph.")
    c.drawString(50, 680, "Second paragraph with updated content.")
    c.drawString(50, 660, "A new third paragraph added.")
    c.drawString(50, 640, "Fourth paragraph with more information.")
    c.save()

    print("PDF files generated successfully (using reportlab)")
    return True


def main():
    """Generate all test fixtures."""
    print("Generating test fixtures...")
    print("-" * 40)

    results = {
        'images': generate_images(),
        'word': generate_word_docs(),
        'excel': generate_excel_files(),
        'powerpoint': generate_powerpoint_files(),
        'pdf': generate_pdf_files(),
    }

    print("-" * 40)
    print("Generation summary:")
    for name, success in results.items():
        status = "OK" if success else "SKIPPED"
        print(f"  {name}: {status}")

    return all(results.values())


if __name__ == '__main__':
    success = main()
    sys.exit(0 if success else 1)
