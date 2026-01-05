import os
from typing import List
from .base import BaseConverter, ConvertedDocument, TextBlock


class PPTXConverter(BaseConverter):
    """Converter for Microsoft PowerPoint PPTX files."""

    @property
    def supported_extensions(self) -> List[str]:
        return ['.pptx', '.ppt']

    def convert(self, file_path: str) -> ConvertedDocument:
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        blocks = []
        full_text = ""

        try:
            from pptx import Presentation
            prs = Presentation(file_path)

            for slide_num, slide in enumerate(prs.slides):
                full_text += f"=== Slide {slide_num + 1} ===\n"
                blocks.append(TextBlock(
                    text=f"=== Slide {slide_num + 1} ===",
                    page=slide_num,
                    x=0,
                    y=0,
                    width=200,
                    height=20,
                    metadata={'type': 'slide_header'}
                ))

                y_offset = 24
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        full_text += text + "\n"

                        left = shape.left / 914400 if shape.left else 0
                        top = shape.top / 914400 if shape.top else y_offset
                        width = shape.width / 914400 if shape.width else 100
                        height = shape.height / 914400 if shape.height else 20

                        blocks.append(TextBlock(
                            text=text,
                            page=slide_num,
                            x=left,
                            y=top,
                            width=width,
                            height=height,
                            metadata={'shape_type': shape.shape_type.name if hasattr(shape.shape_type, 'name') else str(shape.shape_type)}
                        ))
                        y_offset += 20

                    if hasattr(shape, "table"):
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                text = " | ".join(row_text)
                                full_text += text + "\n"
                                blocks.append(TextBlock(
                                    text=text,
                                    page=slide_num,
                                    x=0,
                                    y=y_offset,
                                    width=len(text) * 7,
                                    height=14,
                                    metadata={'type': 'table_row'}
                                ))
                                y_offset += 14

                full_text += "\n"

        except ImportError:
            import subprocess
            try:
                result = subprocess.run(
                    ['pandoc', '-t', 'plain', file_path],
                    capture_output=True,
                    text=True,
                    check=True
                )
                full_text = result.stdout
                lines = full_text.split('\n')
                y_offset = 0
                for line in lines:
                    if line.strip():
                        blocks.append(TextBlock(
                            text=line,
                            page=0,
                            x=0,
                            y=y_offset,
                            width=len(line) * 7,
                            height=12
                        ))
                    y_offset += 12
            except FileNotFoundError:
                raise RuntimeError(
                    "Neither python-pptx nor pandoc is available. "
                    "Install python-pptx: pip install python-pptx"
                )

        return ConvertedDocument(
            blocks=blocks,
            full_text=full_text,
            page_count=len(prs.slides) if 'prs' in locals() else 1,
            source_path=file_path,
            source_type='pptx'
        )
